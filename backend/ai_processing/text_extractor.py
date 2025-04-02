# mrAssignment/ai_processing/text_extractor.py
import logging
import io
import os
import tempfile
from typing import Dict, Optional, Tuple

# PDF processing
import PyPDF2
import pdfplumber

# Document processing
try:
    import docx
except ImportError:
    docx = None

# Presentation processing
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# For detecting MIME types
import mimetypes

logger = logging.getLogger(__name__)

class TextExtractor:
    """
    Utility class for extracting text from various file formats.
    Supports PDF, DOCX, PPTX, and plain text files.
    """
    
    @staticmethod
    def extract_text(file_content: bytes, file_name: str = None) -> Tuple[str, Dict, int]:
        """
        Extract text from file content based on file extension or detected MIME type.
        
        Args:
            file_content (bytes): Binary file content
            file_name (str, optional): Original filename with extension
            
        Returns:
            tuple: (extracted_text, metadata, page_count)
                - extracted_text (str): The extracted text content
                - metadata (dict): Additional metadata about the extraction
                - page_count (int): Number of pages in the document
        """
        if file_name:
            file_extension = os.path.splitext(file_name)[1].lower()
        else:
            # Try to detect MIME type from content
            file_extension = ".txt"  # Default
            
        # Create a file-like object from the byte content
        file_stream = io.BytesIO(file_content)
        
        # Extract based on file extension
        if file_extension in ('.pdf'):
            return TextExtractor._extract_from_pdf(file_stream)
        elif file_extension in ('.docx', '.doc'):
            return TextExtractor._extract_from_docx(file_stream)
        elif file_extension in ('.pptx', '.ppt'):
            return TextExtractor._extract_from_pptx(file_stream)
        elif file_extension in ('.txt', '.md', '.rtf'):
            return TextExtractor._extract_from_text(file_stream)
        else:
            # Default to trying to read as text
            logger.warning(f"Unsupported file extension: {file_extension}. Attempting to extract as plain text.")
            return TextExtractor._extract_from_text(file_stream)
    
    @staticmethod
    def _extract_from_pdf(file_stream: io.BytesIO) -> Tuple[str, Dict, int]:
        """
        Extract text from PDF file using a combination of PyPDF2 and pdfplumber.
        
        Args:
            file_stream (io.BytesIO): PDF file content as BytesIO
            
        Returns:
            tuple: (extracted_text, metadata, page_count)
        """
        full_text = ""
        metadata = {}
        page_count = 0
        
        try:
            # Try with pdfplumber first (better at text extraction)
            with pdfplumber.open(file_stream) as pdf:
                page_count = len(pdf.pages)
                pages_text = []
                
                # Get metadata
                if hasattr(pdf, 'metadata') and pdf.metadata:
                    metadata = {k: str(v) for k, v in pdf.metadata.items()}
                
                # Extract each page
                for page in pdf.pages:
                    page_text = page.extract_text() or ""
                    pages_text.append(page_text)
                
                full_text = "\n\n".join(pages_text)
                
            # If pdfplumber didn't get text, try PyPDF2 as fallback
            if not full_text.strip():
                file_stream.seek(0)  # Reset stream position
                reader = PyPDF2.PdfReader(file_stream)
                page_count = len(reader.pages)
                
                # Get metadata
                if reader.metadata:
                    metadata = {k: str(v) for k, v in reader.metadata.items()}
                
                # Extract each page
                pages_text = []
                for page in reader.pages:
                    page_text = page.extract_text() or ""
                    pages_text.append(page_text)
                
                full_text = "\n\n".join(pages_text)
            
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            full_text = ""
            metadata = {"extraction_error": str(e)}
        
        return full_text, metadata, page_count
    
    @staticmethod
    def _extract_from_docx(file_stream: io.BytesIO) -> Tuple[str, Dict, int]:
        """
        Extract text from DOCX file.
        
        Args:
            file_stream (io.BytesIO): DOCX file content as BytesIO
            
        Returns:
            tuple: (extracted_text, metadata, page_count)
        """
        if not docx:
            logger.error("python-docx package is not installed")
            return "", {"extraction_error": "python-docx not installed"}, 0
        
        try:
            doc = docx.Document(file_stream)
            
            # Extract text
            full_text = "\n\n".join([para.text for para in doc.paragraphs])
            
            # Get metadata
            metadata = {
                "title": doc.core_properties.title or "",
                "author": doc.core_properties.author or "",
                "created": str(doc.core_properties.created) if doc.core_properties.created else "",
                "modified": str(doc.core_properties.modified) if doc.core_properties.modified else ""
            }
            
            # Page count is trickier for DOCX, estimating
            page_count = max(1, len(doc.paragraphs) // 10)  # Rough estimate
            
            return full_text, metadata, page_count
            
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {str(e)}")
            return "", {"extraction_error": str(e)}, 0
            
    @staticmethod
    def _extract_from_pptx(file_stream: io.BytesIO) -> Tuple[str, Dict, int]:
        """
        Extract text from PPTX file.
        
        Args:
            file_stream (io.BytesIO): PPTX file content as BytesIO
            
        Returns:
            tuple: (extracted_text, metadata, page_count)
        """
        if not Presentation:
            logger.error("python-pptx package is not installed")
            return "", {"extraction_error": "python-pptx not installed"}, 0
            
        try:
            presentation = Presentation(file_stream)
            
            # Extract text from each slide
            slides_text = []
            for slide in presentation.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                slides_text.append("\n".join(slide_text))
            
            full_text = "\n\n===== Next Slide =====\n\n".join(slides_text)
            
            # Minimal metadata
            metadata = {
                "slide_count": len(presentation.slides),
            }
            
            return full_text, metadata, len(presentation.slides)
            
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {str(e)}")
            return "", {"extraction_error": str(e)}, 0
    
    @staticmethod
    def _extract_from_text(file_stream: io.BytesIO) -> Tuple[str, Dict, int]:
        """
        Extract text from plain text file.
        
        Args:
            file_stream (io.BytesIO): Text file content as BytesIO
            
        Returns:
            tuple: (extracted_text, metadata, page_count)
        """
        try:
            # Try common encodings
            encodings = ['utf-8', 'latin-1', 'ascii']
            text = None
            
            for encoding in encodings:
                try:
                    file_stream.seek(0)
                    text = file_stream.read().decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                # If all attempts fail, try with errors='replace'
                file_stream.seek(0)
                text = file_stream.read().decode('utf-8', errors='replace')
            
            # Count lines as rough page estimate
            lines = text.count('\n') + 1
            page_count = max(1, lines // 30)  # Rough estimate: ~30 lines per page
            
            return text, {"encoding": encoding}, page_count
            
        except Exception as e:
            logger.error(f"Error extracting text from text file: {str(e)}")
            return "", {"extraction_error": str(e)}, 0


def chunk_text(text: str, chunk_size: int = 1500, overlap: int = 150) -> list:
    """
    Split a text into smaller chunks with specified size and overlap.
    
    Args:
        text (str): The text to chunk
        chunk_size (int): Target size of each chunk in characters
        overlap (int): Number of characters to overlap between chunks
        
    Returns:
        list: List of text chunks
    """
    if not text or not text.strip():
        return []
        
    # Use paragraphs as natural break points
    paragraphs = text.split('\n\n')
    chunks = []
    current_chunk = []
    current_size = 0
    
    for paragraph in paragraphs:
        paragraph_clean = paragraph.strip()
        paragraph_size = len(paragraph_clean)
        
        if paragraph_size == 0:
            continue
            
        # If adding this paragraph exceeds the chunk size and we already have content
        if current_size + paragraph_size > chunk_size and current_size > 0:
            # Save the current chunk
            chunks.append('\n\n'.join(current_chunk))
            
            # Start a new chunk with overlap by keeping some paragraphs
            overlap_size = 0
            overlap_paragraphs = []
            
            # Add paragraphs from the end until we reach desired overlap
            for p in reversed(current_chunk):
                if overlap_size + len(p) <= overlap:
                    overlap_paragraphs.insert(0, p)
                    overlap_size += len(p)
                else:
                    break
                    
            # Start new chunk with overlap paragraphs
            current_chunk = overlap_paragraphs
            current_size = overlap_size
        
        # Add the current paragraph to the chunk
        current_chunk.append(paragraph_clean)
        current_size += paragraph_size
    
    # Don't forget the last chunk
    if current_chunk:
        chunks.append('\n\n'.join(current_chunk))
    
    return chunks